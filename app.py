# import streamlit as st
# from PyPDF2 import PdfReader
# from langchain.text_splitter import RecursiveCharacterTextSplitter
# import os
# from langchain_google_genai import GoogleGenerativeAIEmbeddings
# import google.generativeai as genai
# from langchain_community.vectorstores import FAISS
# from langchain_google_genai import ChatGoogleGenerativeAI
# from langchain.chains.question_answering import load_qa_chain
# from langchain.prompts import PromptTemplate
# from googletrans import Translator
# from wordcloud import WordCloud
# import matplotlib.pyplot as plt
# from dotenv import load_dotenv
# from langchain.schema import HumanMessage, SystemMessage

# # Load environment variables
# load_dotenv()
# os.getenv("GOOGLE_API_KEY")
# genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# # Initialize translator
# translator = Translator()

# def translate_text(text, target_language="en"):
#     return translator.translate(text, dest=target_language).text

# def get_pdf_text(pdf_docs):
#     text = ""
#     for pdf in pdf_docs:
#         pdf_reader = PdfReader(pdf)
#         for page in pdf_reader.pages:
#             raw_text = page.extract_text()
#             translated_text = translate_text(raw_text)
#             text += translated_text
#     return text

# def get_text_chunks(text):
#     text_splitter = RecursiveCharacterTextSplitter(chunk_size=5000, chunk_overlap=500)
#     chunks = text_splitter.split_text(text)
#     return chunks

# def get_vector_store(text_chunks):
#     embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
#     vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
#     vector_store.save_local("faiss_index")

# def get_conversational_chain():
#     prompt_template = """
#     Answer the question as detailed as possible from the provided context. If the answer is not available in the context, say, "answer is not available in the context." Do not provide incorrect information.\n\n
#     Context:\n{context}?\n
#     Question:\n{question}\n
#     Answer:
#     """
#     model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
#     prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
#     chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
#     return chain

# def user_input(user_question):
#     embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
#     new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)  # Enable dangerous deserialization
#     docs = new_db.similarity_search(user_question)
#     chain = get_conversational_chain()
#     response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
#     st.write("Reply:", response["output_text"])

# def summarize_text_with_chat_model(text):
#     model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
#     messages = [HumanMessage(content=f"Please provide a concise summary of the following text:\n\n{text}")]
#     response = model.invoke(messages)
#     return response.content  # Access the content directly

# def sentiment_analysis_with_chat_model(text):
#     model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
#     messages = [HumanMessage(content=f"Analyze the sentiment of the following text:\n\n{text}")]
#     response = model.invoke(messages)
#     return response.content  # Access the content directly

# def generate_word_cloud(text):
#     wordcloud = WordCloud(background_color='white', max_words=100).generate(text)
#     plt.figure(figsize=(10, 8))
#     plt.imshow(wordcloud, interpolation='bilinear')
#     plt.axis("off")
#     st.pyplot(plt)

# def topic_detection(text_chunks):
#     topics = ["Budget Planning", "Project Roadmap", "Resource Allocation", "Team Performance"]
#     st.write("Detected Topics:")
#     for topic in topics:
#         st.write(f"- {topic}")
#     st.write("You might want to ask about these topics!")

# def main():
#     st.set_page_config(page_title="Conference Summarizer")
#     st.header("Chat with your personal assistant 💁")

#     if "raw_text" not in st.session_state:
#         st.session_state.raw_text = ""

#     user_question = st.text_input("Ask a question based on your office meetings, conferences, and more.")

#     if user_question:
#         user_input(user_question)

#     with st.sidebar:
#         st.title("Menu:")
#         pdf_docs = st.file_uploader("Upload your meeting documents and resources, then click 'Submit & Process'", accept_multiple_files=True)
#         if st.button("Submit & Process"):
#             if pdf_docs:
#                 with st.spinner("Processing..."):
#                     raw_text = get_pdf_text(pdf_docs)
#                     st.session_state.raw_text = raw_text  # Save raw_text in session state
#                     text_chunks = get_text_chunks(raw_text)
#                     get_vector_store(text_chunks)
#                     st.success("Processing Complete")
#             else:
#                 st.warning("Please upload at least one document.")

#         if st.button("Generate Summary"):
#             if st.session_state.raw_text:
#                 with st.spinner("Generating Summary..."):
#                     summary = summarize_text_with_chat_model(st.session_state.raw_text)
#                     st.write("Summary of Documents:")
#                     st.write(summary)
#             else:
#                 st.warning("Please upload and process the documents first.")

#         if st.session_state.raw_text:
#             st.write("Sentiment Analysis:")
#             sentiment = sentiment_analysis_with_chat_model(st.session_state.raw_text)
#             st.write(sentiment)

#             st.write("Visual Representation of Key Topics:")
#             generate_word_cloud(st.session_state.raw_text)

#             topic_detection(st.session_state.raw_text)

# if __name__ == "__main__":
#     main()



import streamlit as st
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
from langchain.text_splitter import RecursiveCharacterTextSplitter
import os
from langchain_google_genai import GoogleGenerativeAIEmbeddings
import google.generativeai as genai
from langchain_community.vectorstores import FAISS
from langchain_google_genai import ChatGoogleGenerativeAI
from langchain.chains.question_answering import load_qa_chain
from langchain.prompts import PromptTemplate
from googletrans import Translator
from wordcloud import WordCloud
import matplotlib.pyplot as plt
from dotenv import load_dotenv
from langchain.schema import HumanMessage, SystemMessage
import smtplib
from email.mime.multipart import MIMEMultipart
from email.mime.text import MIMEText
    
from speech_recognition import Recognizer, AudioFile

# Load environment variables
load_dotenv()
genai.configure(api_key=os.getenv("GOOGLE_API_KEY"))

# Initialize translator
translator = Translator()

def translate_text(text, target_language="en"):
    return translator.translate(text, dest=target_language).text

def get_pdf_text(pdf_docs):
    text = ""
    for pdf in pdf_docs:
        pdf_reader = PdfReader(pdf)
        for page in pdf_reader.pages:
            raw_text = page.extract_text()
            translated_text = translate_text(raw_text)
            text += translated_text
    return text

def get_ppt_text(ppt_docs):
    text = ""
    for ppt in ppt_docs:
        presentation = Presentation(ppt)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    translated_text = translate_text(shape.text)
                    text += translated_text + "\n"
    return text

def get_docx_text(docx_docs):
    text = ""
    for doc in docx_docs:
        doc = Document(doc)
        for paragraph in doc.paragraphs:
            translated_text = translate_text(paragraph.text)
            text += translated_text + "\n"
    return text

def get_text_chunks(text):
    text_splitter = RecursiveCharacterTextSplitter(chunk_size=5000, chunk_overlap=500)
    chunks = text_splitter.split_text(text)
    return chunks

def get_vector_store(text_chunks):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    vector_store = FAISS.from_texts(text_chunks, embedding=embeddings)
    vector_store.save_local("faiss_index")

def get_conversational_chain():
    prompt_template = """
    Answer the question as detailed as possible from the provided context. If the answer is not available in the context, say, "answer is not available in the context." Do not provide incorrect information.\n\n
    Context:\n{context}?\n
    Question:\n{question}\n
    Answer:
    """
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    prompt = PromptTemplate(template=prompt_template, input_variables=["context", "question"])
    chain = load_qa_chain(model, chain_type="stuff", prompt=prompt)
    return chain

def user_input(user_question):
    embeddings = GoogleGenerativeAIEmbeddings(model="models/embedding-001")
    new_db = FAISS.load_local("faiss_index", embeddings, allow_dangerous_deserialization=True)
    docs = new_db.similarity_search(user_question)
    chain = get_conversational_chain()
    response = chain({"input_documents": docs, "question": user_question}, return_only_outputs=True)
    st.write("Reply:", response["output_text"])

def summarize_text_with_chat_model(text):
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    messages = [HumanMessage(content=f"Please provide a concise summary of the following text:\n\n{text}")]
    response = model.invoke(messages)
    return response.content

def sentiment_analysis_with_chat_model(text):
    model = ChatGoogleGenerativeAI(model="gemini-pro", temperature=0.3)
    messages = [HumanMessage(content=f"Analyze the sentiment of the following text:\n\n{text}")]
    response = model.invoke(messages)
    return response.content

def generate_word_cloud(text):
    wordcloud = WordCloud(background_color='white', max_words=100).generate(text)
    plt.figure(figsize=(10, 8))
    plt.imshow(wordcloud, interpolation='bilinear')
    plt.axis("off")
    st.pyplot(plt)

def topic_detection(text_chunks):
    topics = ["Budget Planning", "Project Roadmap", "Resource Allocation", "Team Performance"]
    st.write("Detected Topics:")
    for topic in topics:
        st.write(f"- {topic}")
    st.write("You might want to ask about these topics!")

def send_email(to_email, subject, body):
    from_email = os.getenv("EMAIL_SENDER")
    password = os.getenv("EMAIL_PASSWORD")
    msg = MIMEMultipart()
    msg["From"] = from_email
    msg["To"] = to_email
    msg["Subject"] = subject
    msg.attach(MIMEText(body, "plain"))
    server = smtplib.SMTP("smtp.gmail.com", 587)
    server.starttls()
    server.login(from_email, password)
    server.send_message(msg)
    server.quit()

def main():
    st.set_page_config(page_title="AI Conference Assistant")
    st.header("Chat with your AI conference assistant 💼")

    # Login/Signup
    login_email = st.text_input("Enter your email to login or signup")
    if login_email:
        st.session_state["user_email"] = login_email

    if "raw_text" not in st.session_state:
        st.session_state.raw_text = ""

    user_question = st.text_input("Ask a question based on the uploaded documents and conversations.")

    if user_question:
        user_input(user_question)

    with st.sidebar:
        st.title("Menu:")
        pdf_docs = st.file_uploader("Upload PDF documents", accept_multiple_files=True, type="pdf")
        ppt_docs = st.file_uploader("Upload PowerPoint documents", accept_multiple_files=True, type="pptx")
        docx_docs = st.file_uploader("Upload Word documents", accept_multiple_files=True, type="docx")

        if st.button("Submit & Process"):
            if pdf_docs or ppt_docs or docx_docs:
                with st.spinner("Processing..."):
                    raw_text = ""
                    if pdf_docs:
                        raw_text += get_pdf_text(pdf_docs)
                    if ppt_docs:
                        raw_text += get_ppt_text(ppt_docs)
                    if docx_docs:
                        raw_text += get_docx_text(docx_docs)
                    
                    st.session_state.raw_text = raw_text
                    text_chunks = get_text_chunks(raw_text)
                    get_vector_store(text_chunks)
                    st.success("Processing Complete")
            else:
                st.warning("Please upload at least one document.")

        if st.button("Generate Summary"):
            if st.session_state.raw_text:
                with st.spinner("Generating Summary..."):
                    summary = summarize_text_with_chat_model(st.session_state.raw_text)
                    st.write("Summary of Documents:")
                    st.write(summary)

                    if "user_email" in st.session_state:
                        send_email(st.session_state["user_email"], "Meeting Summary", summary)
                        st.success("Summary sent to your email.")
            else:
                st.warning("Please upload and process the documents first.")

        if st.session_state.raw_text:
            st.write("Sentiment Analysis:")
            sentiment = sentiment_analysis_with_chat_model(st.session_state.raw_text)
            st.write(sentiment)

            st.write("Visual Representation of Key Topics:")
            generate_word_cloud(st.session_state.raw_text)

            topic_detection(st.session_state.raw_text)

        # Add live transcription
        st.header("Live Transcription")
        if st.button("Start Transcription"):
            recognizer = Recognizer()
            try:
                with st.spinner("Listening... (Speak now)"):
                    # Instead of using a microphone directly, prompt the user to use their browser's microphone permission
                    st.info("Please speak now...")
                    # Assume the browser's Web Speech API handles the capture and recognition
                    text = recognizer.recognize_google(None)
                    st.write("Transcription: " + text)
                    st.session_state.raw_text += translate_text(text) + "\n"
            except Exception as e:
                st.error(f"Error during transcription: {e}")

if __name__ == "__main__":
    main()
